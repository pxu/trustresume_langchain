"""Generate docs/weekly-submissions/Week8_Demonstration.pptx.

Throwaway doc-tooling, like build_week7_presentation.py. Not part of the app,
not covered by the test suite. This is the *short* deck that frames the Week 8
live demonstration — it bookends the demo rather than replacing it: a title,
one thesis slide, an agenda mapping to the demo storyboard, an architecture
glance, a "live demo" divider, one "watch for this" slide calling out the
honest-failure moment, and a close. The actual walkthrough happens in the
running app (see docs/weekly-submissions/Week8_Demo_Script.md for the
storyboard those beats come from).

Deliberately self-contained (its own slide helpers, no chart imports) so it
doesn't drag in Week 7's report figures — a demo deck is light on charts. The
architecture slide reuses the already-generated docs/diagrams/
01_system_overview.png. Every slide's speaker notes hold the narration;
docs/weekly-submissions/Week8_Presentation_Script.md holds the same text as a
plain read-through transcript for rehearsal.

    .venv/bin/python scripts/build_week8_presentation.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUTPUT_PPTX = ROOT / "docs" / "weekly-submissions" / "Week8_Demonstration.pptx"
OUTPUT_SCRIPT = ROOT / "docs" / "weekly-submissions" / "Week8_Presentation_Script.md"
ARCH_IMAGE = ROOT / "docs" / "diagrams" / "01_system_overview.png"

BLUE = RGBColor(0x4C, 0x72, 0xB0)
ORANGE = RGBColor(0xDD, 0x84, 0x52)
GREEN = RGBColor(0x3E, 0x8E, 0x5A)
DARK = RGBColor(0x2A, 0x2A, 0x2A)
GRAY = RGBColor(0x6B, 0x6B, 0x6B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF7, 0xF9, 0xFC)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.7)
FONT = "Calibri"
FOOTER = "TrustResume | MSAI 699 Capstone — Week 8 Demonstration"

SCRIPT_PARTS: list[tuple[str, str]] = []


def _blank_slide(prs: Presentation):  # type: ignore[no-untyped-def]
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def _add_textbox(slide, left, top, width, height):  # type: ignore[no-untyped-def]
    box = slide.shapes.add_textbox(left, top, width, height)
    box.text_frame.word_wrap = True
    return box


def _style_run(run, *, size: int, color: RGBColor = DARK, bold: bool = False, italic: bool = False) -> None:  # type: ignore[no-untyped-def]
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic


def set_notes(slide, text: str, title: str) -> None:  # type: ignore[no-untyped-def]
    slide.notes_slide.notes_text_frame.text = text
    SCRIPT_PARTS.append((title, text))


def add_background(slide, color: RGBColor) -> None:  # type: ignore[no-untyped-def]
    rect = slide.shapes.add_shape(1, Emu(0), Emu(0), SLIDE_W, SLIDE_H)  # RECTANGLE == 1
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    rect.shadow.inherit = False
    slide.shapes._spTree.remove(rect._element)
    slide.shapes._spTree.insert(2, rect._element)


def add_header(slide, title: str, *, bar_color: RGBColor = BLUE) -> None:  # type: ignore[no-untyped-def]
    bar = slide.shapes.add_shape(1, Emu(0), Emu(0), SLIDE_W, Inches(1.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = bar_color
    bar.line.fill.background()
    bar.shadow.inherit = False
    box = _add_textbox(slide, MARGIN, Inches(0.18), SLIDE_W - 2 * MARGIN, Inches(0.85))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = title
    _style_run(run, size=30, color=WHITE, bold=True)


def add_footer(slide, text: str = FOOTER) -> None:  # type: ignore[no-untyped-def]
    box = _add_textbox(slide, MARGIN, SLIDE_H - Inches(0.45), SLIDE_W - 2 * MARGIN, Inches(0.35))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = text
    _style_run(run, size=11, color=GRAY)


def add_bullets(slide, items: list[str], *, left, top, width, height, size: int = 22, gap: int = 16) -> None:  # type: ignore[no-untyped-def]
    box = _add_textbox(slide, left, top, width, height)
    tf = box.text_frame
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        run = p.add_run()
        run.text = f"•  {item}"
        _style_run(run, size=size, color=DARK)


def content_slide(prs, title, bullets, notes, *, bar_color=BLUE, size=22) -> None:  # type: ignore[no-untyped-def]
    slide = _blank_slide(prs)
    add_background(slide, LIGHT_BG)
    add_header(slide, title, bar_color=bar_color)
    add_bullets(slide, bullets, left=MARGIN, top=Inches(1.6), width=SLIDE_W - 2 * MARGIN, height=Inches(5.3), size=size)
    add_footer(slide)
    set_notes(slide, notes, title)


def image_slide(prs, title, image_path, caption, notes, *, bar_color=BLUE, img_width=Inches(10.8)) -> None:  # type: ignore[no-untyped-def]
    slide = _blank_slide(prs)
    add_background(slide, WHITE)
    add_header(slide, title, bar_color=bar_color)
    left = (SLIDE_W - img_width) / 2
    pic = slide.shapes.add_picture(str(image_path), left, Inches(1.45), width=img_width)
    box = _add_textbox(slide, MARGIN, pic.top + pic.height + Inches(0.1), SLIDE_W - 2 * MARGIN, Inches(0.5))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = caption
    _style_run(run, size=14, color=GRAY, italic=True)
    add_footer(slide)
    set_notes(slide, notes, title)


def divider_slide(prs, kicker, title, notes) -> None:  # type: ignore[no-untyped-def]
    slide = _blank_slide(prs)
    add_background(slide, DARK)
    box0 = _add_textbox(slide, Inches(1.0), Inches(2.6), Inches(11.3), Inches(0.7))
    r0 = box0.text_frame.paragraphs[0].add_run()
    r0.text = kicker
    _style_run(r0, size=22, color=ORANGE, bold=True)
    box = _add_textbox(slide, Inches(1.0), Inches(3.3), Inches(11.3), Inches(1.4))
    r = box.text_frame.paragraphs[0].add_run()
    r.text = title
    _style_run(r, size=46, color=WHITE, bold=True)
    set_notes(slide, notes, kicker)


def build() -> Presentation:  # type: ignore[no-untyped-def]
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # --- 1. Title -------------------------------------------------------
    slide = _blank_slide(prs)
    add_background(slide, BLUE)
    box = _add_textbox(slide, Inches(1.0), Inches(2.2), Inches(11.3), Inches(1.4))
    r = box.text_frame.paragraphs[0].add_run()
    r.text = "TrustResume — Live Demonstration"
    _style_run(r, size=48, color=WHITE, bold=True)
    box2 = _add_textbox(slide, Inches(1.0), Inches(3.7), Inches(11.3), Inches(1.0))
    r2 = box2.text_frame.paragraphs[0].add_run()
    r2.text = "Evidence-based resume generation that would rather fail than lie"
    _style_run(r2, size=22, color=WHITE)
    box3 = _add_textbox(slide, Inches(1.0), Inches(5.7), Inches(11.3), Inches(1.0))
    r3 = box3.text_frame.paragraphs[0].add_run()
    r3.text = "Peng Fei Xu  |  MSAI 699 Capstone — Week 8: Project Demonstration"
    _style_run(r3, size=16, color=RGBColor(0xE3, 0xEA, 0xF7))
    set_notes(
        slide,
        "Hi, I'm Peng Fei Xu. This is the Week 8 demonstration of TrustResume. "
        "Everything you're about to see is the live application running against "
        "a real production model — not slides, not a mock-up. In one sentence: "
        "TrustResume generates a resume grounded in the candidate's own "
        "documents, and then an independent agent checks every claim against "
        "that evidence before any score is shown to anyone.",
        "1. Title",
    )

    # --- 2. Thesis in one slide ----------------------------------------
    content_slide(
        prs,
        "The Idea in One Sentence",
        [
            "AI can write a polished resume in seconds — and just as easily "
            "overstate scope, seniority, or a metric the evidence doesn't support.",
            "TrustResume splits the job in two: a Writer agent drafts from the "
            "candidate's own documents; an independent Trust Harness then checks "
            "every claim against that same evidence.",
            "The writer never grades its own homework — that separation is what "
            "makes a claim checkable, not just plausible.",
            "The design bet you'll see proven live: the system would rather fail "
            "a candidate than lie for them.",
        ],
        "Here's the whole idea before we open the app. Generative AI writing a "
        "convincing resume is a solved problem — the hard part is trust. Nothing "
        "stops a model from quietly stretching the truth to make a draft look "
        "better. So TrustResume splits generation from verification into two "
        "separate agents. The writer drafts from the candidate's real documents; "
        "an independent Trust Harness then checks every single claim against that "
        "evidence. The writer never scores itself. The bet the whole project "
        "rests on is that the system should rather fail a candidate than lie for "
        "them — and I'll show that happening on a real run in a moment.",
        bar_color=ORANGE,
    )

    # --- 3. What you'll see (agenda) -----------------------------------
    content_slide(
        prs,
        "What You'll See in This Demo",
        [
            "Documents — upload a résumé; it's parsed, chunked, and embedded into "
            "the candidate's private, user-scoped evidence store.",
            "Jobs — a job is a persisted entity: extract the posting once, then "
            "reuse that extraction for every generation against it.",
            "Generate — hybrid retrieval → Writer → Trust Harness → ATS scoring, "
            "always rewriting at least once and keeping the best-scoring draft.",
            "The result — real Trust and ATS scores, missing keywords, flagged "
            "claims, and a downloadable PDF/Markdown résumé.",
        ],
        "Quick roadmap for the next few minutes. Four tabs. First, Documents: "
        "I'll upload a résumé and you'll see it ingested — parsed, chunked, "
        "embedded — into an evidence store scoped to just this user. Second, "
        "Jobs: a job here is a saved entity; the posting is extracted once and "
        "reused, so we're not re-parsing it every run. Third, Generate: this "
        "kicks off the pipeline — hybrid vector-plus-keyword retrieval, the "
        "Writer drafting from evidence, the Trust Harness checking each claim, "
        "and ATS keyword scoring. It always rewrites at least once more and "
        "keeps whichever draft actually scored best, not the first that passed. "
        "Finally the result: the real scores, the gaps, any flagged claims, and "
        "an actual exportable document.",
    )

    # --- 4. Architecture at a glance -----------------------------------
    if ARCH_IMAGE.exists():
        image_slide(
            prs,
            "Under the Hood: Six Agents, One Quality Loop",
            ARCH_IMAGE,
            "Streamlit UI → FastAPI → LangGraph orchestrator driving the agents "
            "over a hybrid Chroma + SQLite/FTS5 evidence store.",
            "Before we dive in, one architecture slide so the demo makes sense. "
            "The Streamlit UI talks over HTTP to a FastAPI backend, which drives "
            "a LangGraph orchestrator. The orchestrator sequences the agents: Job "
            "Description and Evidence Retrieval run once, then a loop of Resume "
            "Writer, Trust Harness, and ATS Evaluation, with a cached Candidate "
            "Profile feeding the writer. Retrieval is hybrid — part vector "
            "similarity for paraphrases, part keyword search for exact product "
            "names — over a Chroma vector store plus SQLite full-text search, "
            "everything filtered to one user's own data. That's the isolation "
            "boundary the trust story depends on.",
        )
    else:  # pragma: no cover — diagram is committed; guard only for a fresh checkout
        content_slide(
            prs,
            "Under the Hood: Six Agents, One Quality Loop",
            [
                "Streamlit UI → FastAPI backend → LangGraph orchestrator.",
                "Job Description + Evidence Retrieval run once; Writer ↔ Trust "
                "Harness / ATS run in the quality loop.",
                "Hybrid retrieval (Chroma vectors + SQLite FTS5 keywords), fused "
                "by RRF, always scoped to one user's own data.",
            ],
            "One architecture slide so the demo makes sense. The UI talks to a "
            "FastAPI backend driving a LangGraph orchestrator that sequences the "
            "agents, with hybrid vector-plus-keyword retrieval scoped per user.",
        )

    # --- 5. Live demo divider ------------------------------------------
    divider_slide(
        prs,
        "LIVE DEMO",
        "Switching to the running application",
        "Okay — switching to the live app now. This is the real system against a "
        "real Bedrock model. I'll walk through Documents, Jobs, and a generation, "
        "and then we'll look at what came out.",
    )

    # --- 6. Watch for this ---------------------------------------------
    content_slide(
        prs,
        "The Moment to Watch For",
        [
            "When the evidence doesn't match the job, the Writer does NOT invent "
            "experience to close the gap.",
            "Instead: a high Trust score (every claim is true) alongside a low "
            "ATS score (the résumé honestly doesn't match the posting).",
            "Trust 100 / ATS 0 on the same résumé isn't a bug — it's the entire "
            "product thesis, visible in two numbers.",
            "A run that doesn't pass is shown anyway, labeled honestly, never "
            "hidden or silently upgraded.",
        ],
        "As I run this, here's the one moment I want you to watch for. When a "
        "candidate's evidence doesn't actually match the job, a naive generator "
        "would just invent the missing experience to score higher. This system "
        "won't. You'll see it produce a high Trust score — because every claim it "
        "made is genuinely supported — right next to a low ATS score, because the "
        "résumé honestly doesn't match the posting's keywords. Trust high, ATS "
        "low, on the very same résumé. That's not a failure of the system; that "
        "is the system working exactly as designed. And when a run doesn't clear "
        "the quality gate, we show it anyway, labeled plainly — never hidden, "
        "never quietly upgraded.",
        bar_color=GREEN,
    )

    # --- 7. Close -------------------------------------------------------
    slide = _blank_slide(prs)
    add_background(slide, BLUE)
    box = _add_textbox(slide, Inches(1.0), Inches(2.4), Inches(11.3), Inches(1.2))
    r = box.text_frame.paragraphs[0].add_run()
    r.text = "Grounded, then independently verified"
    _style_run(r, size=40, color=WHITE, bold=True)
    box2 = _add_textbox(slide, Inches(1.0), Inches(3.8), Inches(11.3), Inches(1.6))
    r2 = box2.text_frame.paragraphs[0].add_run()
    r2.text = (
        "Separating the agent that writes from the agent that checks is what "
        "makes this trustworthy rather than merely plausible — and you just saw "
        "it choose an honest result over a flattering, fabricated one."
    )
    _style_run(r2, size=20, color=RGBColor(0xE3, 0xEA, 0xF7))
    box3 = _add_textbox(slide, Inches(1.0), Inches(5.9), Inches(11.3), Inches(0.8))
    r3 = box3.text_frame.paragraphs[0].add_run()
    r3.text = "Thank you — questions welcome."
    _style_run(r3, size=22, color=WHITE, bold=True)
    set_notes(
        slide,
        "To close: separating the agent that writes from the agent that checks "
        "is what makes this trustworthy rather than merely plausible — and you "
        "just watched a real run choose an honest, capped result over a "
        "flattering, fabricated one. That's the whole point. Thank you — happy "
        "to take any questions.",
        "7. Close",
    )

    return prs


def write_script() -> None:
    lines = [
        "# TrustResume: Week 8 Demonstration Speaker Script",
        "",
        "Narration for the slides that bookend the live demo, in slide order. "
        "The demo walkthrough itself is storyboarded in `Week8_Demo_Script.md`; "
        "this covers the framing slides before and after it. The same text is "
        "embedded as speaker notes in `Week8_Demonstration.pptx`.",
        "",
    ]
    for title, text in SCRIPT_PARTS:
        lines += [f"## Slide {title}", "", text, ""]
    OUTPUT_SCRIPT.write_text("\n".join(lines), encoding="utf-8")


def main() -> None:
    prs = build()
    OUTPUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_PPTX)
    write_script()
    print(f"Wrote {OUTPUT_PPTX.relative_to(ROOT)}")
    print(f"Wrote {OUTPUT_SCRIPT.relative_to(ROOT)}")


if __name__ == "__main__":
    main()
