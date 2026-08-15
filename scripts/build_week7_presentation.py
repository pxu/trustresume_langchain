"""Generate docs/weekly-submissions/Week7_Final_Presentation.pptx.

Throwaway doc-tooling, like build_week7_report.py. Not part of the app, not
covered by the test suite. A stakeholder-facing deck (not the academic
register of the written report): the problem, the solution, how it works,
what the numbers show, and what's next. Reuses the same two chart figures as
the Week 7 report (imported from build_week7_report.py) so the whole
deliverable set, report and slides, reads as one consistent story. Every
slide's speaker notes hold the full narration; docs/weekly-submissions/
Week7_Presentation_Script.md holds the same text as a plain read-through
transcript for rehearsal.

    .venv/bin/python scripts/build_week7_presentation.py
"""

from __future__ import annotations

import sys
import tempfile
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))

from build_week7_report import make_pipeline_diagram, make_results_chart  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.text import PP_ALIGN  # noqa: E402
from pptx.util import Emu, Inches, Pt  # noqa: E402

ROOT = Path(__file__).resolve().parents[1]
OUTPUT_PPTX = ROOT / "docs" / "weekly-submissions" / "Week7_Final_Presentation.pptx"
OUTPUT_SCRIPT = ROOT / "docs" / "weekly-submissions" / "Week7_Presentation_Script.md"

BLUE = RGBColor(0x4C, 0x72, 0xB0)
ORANGE = RGBColor(0xDD, 0x84, 0x52)
DARK = RGBColor(0x2A, 0x2A, 0x2A)
GRAY = RGBColor(0x6B, 0x6B, 0x6B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF7, 0xF9, 0xFC)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.7)
FONT = "Calibri"

SCRIPT_PARTS: list[tuple[str, str]] = []


def _blank_slide(prs: Presentation):  # type: ignore[no-untyped-def]
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def _add_textbox(slide, left, top, width, height):  # type: ignore[no-untyped-def]
    box = slide.shapes.add_textbox(left, top, width, height)
    box.text_frame.word_wrap = True
    return box


def _style_run(run, *, size: int, color: RGBColor = DARK, bold: bool = False, italic: bool = False) -> None:  # type: ignore[no-untyped-def]
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic


def set_notes(slide, text: str, title: str) -> None:  # type: ignore[no-untyped-def]
    slide.notes_slide.notes_text_frame.text = text
    SCRIPT_PARTS.append((title, text))


def add_background(slide, color: RGBColor) -> None:  # type: ignore[no-untyped-def]
    rect = slide.shapes.add_shape(1, Emu(0), Emu(0), SLIDE_W, SLIDE_H)  # MSO_SHAPE.RECTANGLE == 1
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    rect.shadow.inherit = False
    slide.shapes._spTree.remove(rect._element)
    slide.shapes._spTree.insert(2, rect._element)


def add_header(slide, title: str, *, bar_color: RGBColor = BLUE) -> None:  # type: ignore[no-untyped-def]
    bar = slide.shapes.add_shape(1, Emu(0), Emu(0), SLIDE_W, Inches(1.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = bar_color
    bar.line.fill.background()
    bar.shadow.inherit = False
    box = _add_textbox(slide, MARGIN, Inches(0.18), SLIDE_W - 2 * MARGIN, Inches(0.85))
    p = box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    _style_run(run, size=30, color=WHITE, bold=True)


def add_footer(slide, text: str) -> None:  # type: ignore[no-untyped-def]
    box = _add_textbox(slide, MARGIN, SLIDE_H - Inches(0.45), SLIDE_W - 2 * MARGIN, Inches(0.35))
    p = box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    _style_run(run, size=11, color=GRAY)


def add_bullets(slide, items: list[str], *, left, top, width, height, size: int = 20, gap: int = 14) -> None:  # type: ignore[no-untyped-def]
    box = _add_textbox(slide, left, top, width, height)
    tf = box.text_frame
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        run = p.add_run()
        run.text = f"•  {item}"
        _style_run(run, size=size, color=DARK)


def content_slide(prs: Presentation, title: str, bullets: list[str], notes: str, *, bar_color: RGBColor = BLUE, size: int = 20) -> None:  # type: ignore[no-untyped-def]
    slide = _blank_slide(prs)
    add_background(slide, LIGHT_BG)
    add_header(slide, title, bar_color=bar_color)
    add_bullets(slide, bullets, left=MARGIN, top=Inches(1.55), width=SLIDE_W - 2 * MARGIN, height=Inches(5.3), size=size)
    add_footer(slide, "TrustResume | MSAI 699 Capstone Final Presentation")
    set_notes(slide, notes, title)


def image_slide(prs: Presentation, title: str, image_path: Path, caption: str, notes: str, *, bar_color: RGBColor = BLUE, img_width: Emu = Inches(10.5)) -> None:  # type: ignore[no-untyped-def]
    slide = _blank_slide(prs)
    add_background(slide, WHITE)
    add_header(slide, title, bar_color=bar_color)
    left = (SLIDE_W - img_width) / 2
    pic = slide.shapes.add_picture(str(image_path), left, Inches(1.45), width=img_width)
    cap_top = pic.top + pic.height + Inches(0.1)
    box = _add_textbox(slide, MARGIN, cap_top, SLIDE_W - 2 * MARGIN, Inches(0.5))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = caption
    _style_run(run, size=14, color=GRAY, italic=True)
    add_footer(slide, "TrustResume | MSAI 699 Capstone Final Presentation")
    set_notes(slide, notes, title)


def table_slide(prs: Presentation, title: str, headers: list[str], rows: list[list[str]], notes: str, *, bar_color: RGBColor = BLUE) -> None:  # type: ignore[no-untyped-def]
    slide = _blank_slide(prs)
    add_background(slide, LIGHT_BG)
    add_header(slide, title, bar_color=bar_color)
    n_rows, n_cols = 1 + len(rows), len(headers)
    table_shape = slide.shapes.add_table(
        n_rows, n_cols, MARGIN, Inches(1.9), SLIDE_W - 2 * MARGIN, Inches(0.55 * n_rows)
    )
    table = table_shape.table
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                _style_run(r, size=16, color=WHITE, bold=True)
    for ri, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(ri, c)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if ri % 2 else LIGHT_BG
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    _style_run(r, size=15, color=DARK)
    add_footer(slide, "TrustResume | MSAI 699 Capstone Final Presentation")
    set_notes(slide, notes, title)


def build(tmpdir: Path) -> Presentation:  # type: ignore[no-untyped-def]
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # --- 1. Title slide -------------------------------------------------
    slide = _blank_slide(prs)
    add_background(slide, BLUE)
    box = _add_textbox(slide, Inches(1.0), Inches(2.3), Inches(11.3), Inches(1.6))
    p = box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "TrustResume"
    _style_run(run, size=54, color=WHITE, bold=True)
    box2 = _add_textbox(slide, Inches(1.0), Inches(3.7), Inches(11.3), Inches(1.0))
    p2 = box2.text_frame.paragraphs[0]
    run2 = p2.add_run()
    run2.text = "Evidence-Based Resume Generation Using RAG, Multi-Agent AI, and Trust Verification"
    _style_run(run2, size=22, color=WHITE)
    box3 = _add_textbox(slide, Inches(1.0), Inches(5.6), Inches(11.3), Inches(1.0))
    p3 = box3.text_frame.paragraphs[0]
    run3 = p3.add_run()
    run3.text = "Peng Fei Xu  |  MSAI 699 Capstone Final Presentation  |  University of the Cumberlands"
    _style_run(run3, size=16, color=RGBColor(0xE3, 0xEA, 0xF7))
    set_notes(
        slide,
        "Good [morning/afternoon] everyone, thank you for the time. I'm "
        "Peng Fei Xu, presenting TrustResume, an evidence-based resume "
        "generation system built on retrieval-augmented generation, "
        "multi-agent AI, and an independent trust-verification layer. This "
        "is the final presentation for my MSAI 699 capstone. I'll walk "
        "through the problem, how the system works, what the data shows, "
        "and where it goes from here.",
        "1. Title",
    )

    # --- 2. Problem statement --------------------------------------------
    content_slide(
        prs,
        "The Problem: AI Can Write a Resume. Should You Trust It?",
        [
            "Generative AI makes it trivial to produce a polished, "
            "professional-sounding resume in seconds.",
            "But nothing stops it from quietly overstating scope, seniority, "
            "or a metric the candidate's real background doesn't support.",
            "That risk lands on everyone: candidates lose credibility, "
            "employers make bad hiring decisions, and every legitimate use "
            "of AI in hiring gets tainted by the ones that lie.",
            "The open question this project answers: can an AI system "
            "generate a strong resume and prove every claim in it is true?",
        ],
        "Let's start with the problem. Generative AI can write a "
        "convincing resume in seconds. That part is solved. The problem "
        "is nothing stops it from quietly stretching the truth: bumping "
        "'2 years' to 'senior', turning a real 64 percent improvement into "
        "a rounder 90 percent. That's the easiest way for a model to make "
        "a draft look better, so without a check, it happens by default. "
        "And the cost isn't abstract: candidates lose credibility, "
        "employers make bad hiring decisions, and every legitimate use of "
        "AI in hiring gets painted with the same brush. So the question "
        "driving this project: can we build a system that generates a "
        "strong resume and can prove, claim by claim, that it's true?",
        bar_color=ORANGE,
    )

    # --- 3. The solution ---------------------------------------------------
    content_slide(
        prs,
        "The Solution: Generate, Then Independently Verify",
        [
            "TrustResume retrieves a candidate's own documents (resumes, "
            "project notes, performance reviews) as grounding evidence. "
            "This is retrieval-augmented generation (RAG).",
            "A writer agent drafts the resume from that evidence, but a "
            "second, independent agent then checks every claim against it.",
            "A resume only ships if it clears both a Trust gate and an "
            "ATS keyword-coverage gate; if it doesn't, the system rewrites "
            "and tries again, up to three times.",
            "The core design bet: the system would rather fail a candidate "
            "than lie for them. We'll see that proven with a real example "
            "later.",
        ],
        "The solution is architectural. First, retrieval-augmented "
        "generation: instead of inventing a resume, we retrieve the "
        "candidate's own documents and ground the draft in that "
        "evidence. Second, the key idea: generation and verification are "
        "separate agents. The writer never grades its own homework. An "
        "independent Trust Harness checks every claim against the "
        "evidence, and a resume only ships once it clears a Trust gate "
        "and an ATS gate; otherwise the system rewrites, up to three "
        "times. The bet this project is built on: it should rather fail "
        "a candidate than lie for them. I'll show a real example later.",
    )

    # --- 4. Architecture diagram -------------------------------------------
    fig1 = make_pipeline_diagram(tmpdir)
    image_slide(
        prs,
        "How It Works: Six Agents, One Quality Loop",
        fig1,
        "Job Description → Evidence Retrieval → Resume Writer ↔ Trust Harness / ATS Evaluation, with a cached Candidate Profile feeding the writer.",
        "Here's the architecture. Job Description and Evidence Retrieval "
        "each run once. Retrieval uses a hybrid search, part vector "
        "similarity, part keyword, so it catches both a paraphrase and an "
        "exact product name. Those feed a loop: the Resume Writer drafts, "
        "the Trust Harness scores every claim as supported, partially "
        "supported, or unsupported, and ATS Evaluation scores keyword "
        "coverage. If both clear the bar, we're done; if not, the system "
        "builds specific feedback from what failed and tries again, up to "
        "three rewrites. A sixth agent, a cached Candidate Profile, feeds "
        "the writer without being recomputed every time. This all runs as "
        "a graph, LangGraph, which makes the loop and the stopping "
        "condition explicit and testable, not buried in ad hoc code.",
    )

    # --- 5. Tech stack / design decisions ----------------------------------
    content_slide(
        prs,
        "Built for Reliability, Not Just a Demo",
        [
            "LangChain + LangGraph orchestrate every agent and the rewrite "
            "loop; ChromaDB + SQLite/FTS5 power hybrid retrieval.",
            "Works with AWS Bedrock, OpenAI, or Google: the LLM provider "
            "is a configuration choice, not a code change.",
            "Sampling temperature is pinned to zero: the Trust score is a "
            "deterministic, code-computed number, not a re-roll of the "
            "dice.",
            "Models are tiered by role: a cheap model fills a form, a "
            "stronger one only where it changes what a candidate reads or "
            "what gets verified.",
            "Every request is scoped to a specific user's own data: the "
            "isolation boundary this whole system depends on.",
        ],
        "A few engineering choices matter here because they make this "
        "deployable, not just a proof of concept. It's provider-agnostic: "
        "it runs on AWS Bedrock, OpenAI, or Google Gemini, chosen by "
        "configuration, not by rewriting code, so there's no vendor "
        "lock-in. Sampling temperature is pinned to zero, because the "
        "headline promise is a deterministic, code-computed trust score. "
        "If the same resume could score differently on two runs, that "
        "promise wouldn't hold. Models are tiered by role: a cheap model "
        "handles simple extraction, a stronger one is reserved for "
        "writing and for the verifier, the one place skimping actually "
        "hurts. And every request is scoped to one user's own data: the "
        "isolation boundary the whole trust story depends on.",
    )

    # --- 6. Evaluation methodology ------------------------------------------
    content_slide(
        prs,
        "How We Know It Works: Measuring, Not Assuming",
        [
            "An automated test suite (474 tests, 99.4% coverage) checks "
            "correctness on every change, offline, with no live model "
            "needed.",
            "A separate offline evaluation harness checks quality: does "
            "retrieval actually find the right evidence, and does the Trust "
            "Harness actually classify claims correctly?",
            "That second check matters because the runtime Trust score is "
            "computed from whatever the Trust Harness reports. A harness "
            "that rubber-stamps everything would score perfectly and fail "
            "invisibly.",
            "Both suites run against hand-labeled ground truth, and the "
            "Trust Harness suite is scored against a real, production "
            "model, not a stand-in.",
        ],
        "How do we know this works, rather than just hoping it does? Two "
        "layers. First, a 474-test automated suite at 99.4 percent "
        "coverage, checking correctness offline on every change. But "
        "correctness isn't quality, so there's a second, separate "
        "evaluation harness scoring the system against labeled ground "
        "truth: does retrieval surface the right evidence, does the Trust "
        "Harness classify claims correctly? That second question matters "
        "most, because the score a user sees comes straight from whatever "
        "the Harness reports. A verifier that marks everything supported "
        "would score perfectly and fail invisibly. So we score the "
        "verifier itself, against labeled examples, on a real model.",
    )

    # --- 7. Results: retrieval + calibration chart --------------------------
    fig2 = make_results_chart(tmpdir)
    image_slide(
        prs,
        "The Numbers: Retrieval Finds It, Verification Catches It",
        fig2,
        "Left: Trust Harness accuracy before/after a calibration fix. Right: Trust/ATS scores from four real generations against a production model.",
        "Here's what the data shows. On the left: retrieval finds the "
        "right evidence 100 percent of the time in our labeled test set. "
        "The Trust Harness chart is more interesting: the first version "
        "of the verifier's instructions only scored 50 percent accuracy, "
        "because it said 'be strict' without defining what that meant, so "
        "it downgraded almost everything by one notch. Explicit "
        "definitions pushed accuracy to over 83 percent, and the "
        "dangerous error, something false marked true, stayed at zero the "
        "whole time. On the right: four real runs against a production "
        "Bedrock model. Three of four fail the quality gate, on purpose. "
        "A gate that never says no isn't checking anything.",
    )

    # --- 8. Results: sample scenarios table ---------------------------------
    table_slide(
        prs,
        "The Proof: It Would Rather Fail Than Lie",
        ["Scenario", "Result", "Trust", "ATS", "What happened"],
        [
            ["1. Strong match", "PASS", "94/90", "90/85", "Real overlap; passed on the first draft"],
            ["2. Partial match", "FAIL", "94/90", "53/85", "Honest gaps: evidence has no domain specifics to invent"],
            ["3. Inflation pressure", "FAIL", "92/90", "50/85", "Writer drifted toward overstatement; harness caught it"],
            ["4. Wrong domain", "FAIL", "100/90", "0/85", "Writer refused to invent experience; every claim true"],
        ],
        "This is the slide I most want you to remember. Same candidate, "
        "four job postings, run against a real production model. "
        "Scenario one is a genuine match: it passes cleanly. Scenario four "
        "matters most: an iOS posting for a backend candidate with zero "
        "iOS experience. The writer refused to invent any. Trust score: "
        "a perfect 100, because every claim was true. ATS score: zero, "
        "because the resume matches no required keyword. Trust 100, ATS "
        "0, same résumé. That's not a bug. That's the entire product "
        "thesis, visible in two numbers. And scenario three shows the "
        "harness earning its keep under pressure: asked for scope beyond "
        "the evidence, the writer drifted toward exaggeration, and "
        "verification caught every one of those claims before anyone saw "
        "them.",
    )

    # --- 9. Quality & cost --------------------------------------------------
    content_slide(
        prs,
        "Quality Engineering Behind the Scenes",
        [
            "474 automated tests, 99.4% coverage, enforced on every code "
            "change via continuous integration.",
            "Every generation is measured: tokens, latency, and dollar "
            "cost, broken down per step, so there's no more guessing what "
            "a run costs.",
            "Real-world cost: roughly $0.32 to $1.07 per résumé, across "
            "four to ten model calls depending on how many rewrites it "
            "takes.",
            "Several of the most important bugs only showed up when the "
            "real system ran end to end, not in a unit test, which is "
            "why we kept real, live testing in the process, not just fast "
            "offline checks.",
        ],
        "A quick word on the engineering behind this, because it's what "
        "makes the last slide's numbers trustworthy rather than lucky. "
        "474 automated tests, over 99 percent coverage, enforced on every "
        "change through continuous integration. Every generation is also "
        "instrumented: tokens, time, and dollar cost per step, which in "
        "practice runs about 32 cents to just over a dollar per résumé. "
        "And I'll be candid: several of the most important bugs, like a "
        "real model's output crashing the PDF export, or a cost-reporting "
        "gap that silently hid Bedrock's real cost, only showed up running "
        "the actual system against a real model. No unit test caught "
        "them, which is why we kept real end-to-end testing in the loop.",
    )

    # --- 10. Ethics ----------------------------------------------------------
    content_slide(
        prs,
        "Responsible by Design",
        [
            "Security: every piece of external text, such as a job "
            "posting or an uploaded document, is clearly marked as data, "
            "not instructions, defending against prompt injection.",
            "Fairness: the same verification that catches fabrication also "
            "prevents the system from rewarding whichever candidate's "
            "background happens to embellish best.",
            "Privacy: every record is isolated per user and filtered "
            "server-side on every search.",
            "Known gap: the deployed API identifies a caller from a "
            "self-reported header. Real authentication is on the roadmap, "
            "not yet built.",
        ],
        "A few words on responsible design, since this touches real "
        "people's careers. Security: any external text, such as a "
        "posting or an uploaded document, is tagged as data, never an "
        "instruction, defending against prompt injection. Fairness: the "
        "same verification that catches fabrication also protects "
        "against a subtler bias: an unchecked generator would favor "
        "whichever candidate's background happens to embellish most "
        "convincingly. Privacy: every record and search is scoped to one "
        "user's own data. And I want to be upfront: the deployed API "
        "currently identifies a caller from a self-reported header. "
        "That's enough to prove isolation works, but not real "
        "authentication. That's the "
        "clearest gap before a real multi-tenant deployment.",
    )

    # --- 11. Future work ------------------------------------------------------
    content_slide(
        prs,
        "What's Next",
        [
            "Real authentication and rate limiting for multi-tenant "
            "production use.",
            "A relevance threshold on retrieval, so a question with no "
            "good answer returns less, not a forced full set of results.",
            "A larger labeled dataset targeting the two claim types the "
            "verifier still occasionally misjudges.",
            "Broader language support for exported résumés, and a clear "
            "path to scale past a single server as demand grows.",
        ],
        "So what's next? First, real authentication and rate limiting, "
        "needed before this serves multiple real customers safely. "
        "Second, a relevance threshold on retrieval, so a question with "
        "no good answer returns fewer results instead of forcing a full "
        "set. Third, growing the labeled dataset around the two specific "
        "claim types the verifier still occasionally misjudges. And "
        "fourth, broader language support on export, plus a clear path "
        "to scale beyond a single server. None of this is open research. "
        "It's concrete, scoped engineering work.",
    )

    # --- 12. Conclusion -------------------------------------------------------
    content_slide(
        prs,
        "Key Takeaways",
        [
            "AI can generate content and independently verify it: "
            "generation and verification don't have to be the same step "
            "trusting itself.",
            "We proved it empirically: a real production model, given a "
            "job it couldn't honestly match, chose a perfect trust score "
            "over a fabricated resume.",
            "Quality is measured, not assumed: an evaluation harness "
            "caught a real flaw in our own verifier, and we fixed it "
            "using the same discipline.",
            "This pattern, grounded generation plus independent "
            "verification, generalizes well beyond résumés, to any "
            "AI-generated content where truthfulness matters.",
        ],
        "To close, three takeaways. One: generation and verification "
        "don't have to be the same step trusting itself. Separating "
        "them is what makes an AI system's claims checkable, not just "
        "plausible. Two: we proved it, not just claimed it. A real "
        "model, faced with a job it couldn't honestly match, chose a "
        "perfect trust score over inventing experience, and our own "
        "evaluation harness caught and let us fix a real flaw in the "
        "verifier before it reached a user. Three: this pattern isn't "
        "specific to résumés. Any AI system generating content where "
        "truthfulness matters can ground it in real evidence and check "
        "it independently before anyone sees it. Thank you, and happy to "
        "take questions.",
        bar_color=ORANGE,
    )

    # --- 13. Thank you / Q&A ---------------------------------------------------
    slide = _blank_slide(prs)
    add_background(slide, BLUE)
    box = _add_textbox(slide, Inches(1.0), Inches(2.9), Inches(11.3), Inches(1.2))
    p = box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Thank You"
    _style_run(run, size=48, color=WHITE, bold=True)
    box2 = _add_textbox(slide, Inches(1.0), Inches(4.1), Inches(11.3), Inches(0.8))
    p2 = box2.text_frame.paragraphs[0]
    run2 = p2.add_run()
    run2.text = "Questions?"
    _style_run(run2, size=24, color=RGBColor(0xE3, 0xEA, 0xF7))
    set_notes(
        slide,
        "Thank you again for your time. I'd be glad to take any "
        "questions.",
        "13. Thank You",
    )

    return prs


def write_script() -> None:
    lines = [
        "# TrustResume: Final Presentation Speaker Script",
        "",
        "Full narration, in slide order. Read-through pace (~140 wpm) targets "
        "roughly 7-8 minutes, inside the assignment's 5-10 minute window.",
        "The same text is embedded as speaker notes in "
        "Week7_Final_Presentation.pptx. Use whichever is more convenient "
        "while recording.",
        "",
    ]
    for title, text in SCRIPT_PARTS:
        lines.append(f"## Slide {title}")
        lines.append("")
        lines.append(text)
        lines.append("")
    OUTPUT_SCRIPT.write_text("\n".join(lines), encoding="utf-8")


def main() -> None:
    with tempfile.TemporaryDirectory() as tmp:
        prs = build(Path(tmp))
        OUTPUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
        prs.save(OUTPUT_PPTX)
    write_script()
    print(f"Wrote {OUTPUT_PPTX.relative_to(ROOT)}")
    print(f"Wrote {OUTPUT_SCRIPT.relative_to(ROOT)}")


if __name__ == "__main__":
    main()
